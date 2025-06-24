import os
import re
import cv2
import subprocess
import numpy as np
import pytesseract
from pathlib import Path
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from django.conf import settings
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pdf2image import convert_from_path
from collections import defaultdict

OUTPUT_DIR = os.path.join(settings.BASE_DIR, "media", "extracted_images")

def add_logo(slide, slide_width, slide_height):
    # Add logo to top-left corner
    logo_path = os.path.join(settings.BASE_DIR, "docs", "mcq_logo.png")
    bg_logo_path = os.path.join(settings.BASE_DIR, "docs", "bg_logo.png")
    logo_left = Inches(0.2)
    logo_top = Inches(0.2)
    logo_width = Inches(1.0)  # Adjust as needed
    slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width)

    img_width = Inches(3.5)
    img_height = Inches(3)

    # Center position
    left = (slide_width - img_width) / 2
    top = (slide_height - img_height) / 2

    # Add image
    slide.shapes.add_picture(bg_logo_path, left, top, width=img_width, height=img_height)


def parse_word_document(doc_path):
    """Parse the Word document and extract questions with their arrangements and directions"""
    doc = Document(doc_path)
    
    content_blocks = []
    current_direction = None
    pending_arrangement = None  # Store arrangement for next question
    
    # Track if we're inside an arrangement block
    in_arrangement = False
    
    for para in doc.paragraphs:
        text = para.text.strip()
        
        if not text:
            continue
            
        # Check if it's a direction
        if "Directions for questions" in text or "DIRECTIONS:" in text:
            # Start new direction
            current_direction = text
            
        # Check if it's an arrangement line
        elif "arrangement" in text.lower() or \
        "final" in text.lower() or \
        "follows" in text.lower():
            in_arrangement = True
            pending_arrangement = text + '\n'
            
        elif in_arrangement:
            # Check for various arrangement patterns
            if (text.startswith('[') or  # Underlined arrangement
                '>' in text[:5] or  # Comparison arrangement
                re.match(r'^[A-Z]\s+[A-Z]', text) or  # Letter arrangement
                re.match(r'^[A-Z]\s+\>', text) or  # Ranking arrangement
                '_' in text or  # Underlined format
                text.startswith('**') or  # Bold text in markdown format
                re.match(r'^[A-Z]+\s+[A-Z]+', text) or  # Multiple letter arrangement
                any(c in text for c in ['→', '↑', '↓', '←'])):  # Directional arrows
                
                # Clean up markdown bold formatting if present
                clean_text = text.replace('**', '').strip()
                if pending_arrangement:
                    pending_arrangement += clean_text + '\n'
                else:
                    pending_arrangement = clean_text + '\n'
                    
                # Check if this is likely the end of arrangement
                if not text.endswith(','):
                    in_arrangement = False
                    
        # Check if it's a question number - improved regex to handle bold markers
        if re.match(r'^\d{1,2}+\.', text):
            # Start new question with pending arrangement
            content_blocks.append({
                'type': 'question',
                'direction': current_direction,
                'arrangement': pending_arrangement,
                'content': text,
                'options': []
            })
            pending_arrangement = None  # Reset for next question
            
        # If we have a current question, check if this is an option
        elif content_blocks and content_blocks[-1]['type'] == 'question':
            # Check if this line is an option (starts with number in parentheses)
            if re.match(r'^\(\d+\)', text) or re.match(r'^\\\(\d+\\\)', text):
                content_blocks[-1]['options'].append(text)
            # Otherwise, append to question text if it's not a table or separator
            elif not (text.startswith('-----') or text.startswith('===') or 
                    '|' in text and len(text.split('|')) > 2 or text in pending_arrangement):
                content_blocks[-1]['content'] += '\n' + text
    return content_blocks


# Image extraction functions from test.py
def convert_docx_to_pdf(docx_path, output_dir):
    """Converts a DOCX file to PDF using LibreOffice"""
    os.makedirs(output_dir, exist_ok=True)

    try:
        subprocess.run([
            "soffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_dir,
            docx_path
        ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except FileNotFoundError:
        raise EnvironmentError("LibreOffice (`soffice`) not found. Please install it.")
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"LibreOffice failed to convert the file: {e}")

    pdf_name = Path(docx_path).stem + ".pdf"
    pdf_path = os.path.join(output_dir, pdf_name)

    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"Expected PDF not found at {pdf_path}")

    return pdf_path


def convert_pdf_to_images(pdf_path, image_dir):
    """Convert PDF pages to high-resolution images"""
    os.makedirs(image_dir, exist_ok=True)
    images = convert_from_path(pdf_path, dpi=300)
    page_paths = []

    for i, img in enumerate(images):
        img_path = os.path.join(image_dir, f"page_{i+1}.png")
        img.save(img_path, "PNG")
        page_paths.append(img_path)
    
    return page_paths


def detect_question_regions_enhanced(page_path, known_questions=None):
    """Detect question numbers and their positions using OCR with enhanced preprocessing"""
    img = cv2.imread(page_path)
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    
    # Enhanced preprocessing
    denoised = cv2.bilateralFilter(gray, 9, 75, 75)
    clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8,8))
    enhanced = clahe.apply(denoised)
    
    # Multiple threshold attempts
    question_regions = []
    
    # Try different preprocessing methods
    preprocessing_methods = [
        lambda x: cv2.threshold(x, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1],
        lambda x: cv2.adaptiveThreshold(x, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2),
        lambda x: x  # Use enhanced image as-is
    ]
    
    for preprocess in preprocessing_methods:
        processed = preprocess(enhanced)
        
        try:
            # Get OCR data with custom config
            custom_config = r'--oem 3 --psm 6'
            data = pytesseract.image_to_data(processed, config=custom_config, output_type=pytesseract.Output.DICT)
            
            # Process OCR results
            for i in range(len(data['text'])):
                text = str(data['text'][i]).strip()
                
                # Multiple patterns for flexibility
                patterns = [
                    r'^(\d{1,2})\.$',           # "14."
                    r'^(\d{1,2})\.\s*$',        # "14. "
                ]
                
                question_num = None
                for pattern in patterns:
                    match = re.match(pattern, text)
                    if match:
                        question_num = int(match.group(1))
                        break
                
                # Also check if this might be part of a question (e.g., separated "14" and ".")
                if not question_num and text.isdigit() and 1 <= int(text) <= 30:
                    # Check if next text element is a period
                    if i + 1 < len(data['text']) and data['text'][i + 1].strip() in ['.', '．']:
                        question_num = int(text)
                
                if question_num and data['conf'][i] > 30:  # Confidence threshold
                    x, y, w, h = data['left'][i], data['top'][i], data['width'][i], data['height'][i]
                    
                    if w > 5 and h > 5:  # Minimum size
                        question_regions.append({
                            'number': question_num,
                            'x': x,
                            'y': y,
                            'width': w,
                            'height': h,
                            'confidence': data['conf'][i]
                        })
            
            # If we found questions with this method, stop trying others
            if len(question_regions) >= 5:  # Expect at least 5 questions per page
                break
                
        except Exception as e:
            continue
    
    # Deduplicate by keeping highest confidence for each question
    unique_questions = {}
    for q in question_regions:
        num = q['number']
        if num not in unique_questions or q['confidence'] > unique_questions[num]['confidence']:
            unique_questions[num] = q
    
    question_regions = list(unique_questions.values())
    
    # Sort by position
    question_regions.sort(key=lambda q: (q['y'], q['x']))
    
    return question_regions


def is_valid_diagram(contour, w, h, area):
    """Check if a contour represents a valid diagram (not a line or text)"""
    # Filter out very small areas
    if area < 5000:  # Increased threshold to avoid small elements
        return False
    
    # Calculate aspect ratio
    aspect_ratio = float(w) / h if h > 0 else float('inf')
    
    # Filter out extreme aspect ratios (lines)
    if aspect_ratio > 10 or aspect_ratio < 0.1:
        return False
    
    # Filter out very thin shapes (likely lines or underlines)
    if w < 20 or h < 20:
        return False
    
    # Calculate solidity (ratio of contour area to convex hull area)
    hull = cv2.convexHull(contour)
    hull_area = cv2.contourArea(hull)
    solidity = float(area) / hull_area if hull_area > 0 else 0
    
    # Filter out very low solidity shapes (likely text or noise)
    if solidity < 0.3:
        return False
    
    return True


def find_associated_question(img_y, img_x, img_h, question_regions, page_width):
    """Find which question a diagram belongs to based on position"""
    if not question_regions:
        return None
    
    # Center position of the image
    img_center_y = img_y + img_h // 2
    
    # Determine column (left or right)
    is_right_column = img_x > page_width / 2
    
    # Filter questions by column
    column_questions = []
    for q in question_regions:
        q_is_right = q['x'] > page_width / 2
        if is_right_column == q_is_right:
            column_questions.append(q)
    
    if not column_questions:
        # If no questions in the same column, use all questions
        column_questions = question_regions
    
    # Find the closest question that is BELOW the image
    closest_question = None
    min_distance = float('inf')
    
    for q in column_questions:
        # Consider questions that are BELOW the image center
        if q['y'] > img_center_y:
            distance = q['y'] - img_center_y
            if distance < min_distance:
                min_distance = distance
                closest_question = q['number']
    
    # If no question found below, try finding the closest question above
    if not closest_question:
        for q in column_questions:
            if q['y'] <= img_center_y:
                distance = img_center_y - q['y']
                if distance < min_distance:
                    min_distance = distance
                    closest_question = q['number']
    
    return closest_question


def extract_diagrams_from_pages(page_paths, output_dir, known_questions=None):
    """Extract diagrams from pages and associate with questions"""
    os.makedirs(output_dir, exist_ok=True)
    results = defaultdict(list)  # Question -> list of image paths
    
    for page_idx, page_path in enumerate(page_paths):
        img = cv2.imread(page_path)
        page_height, page_width = img.shape[:2]
        
        # Detect question regions on this page using OCR
        question_regions = detect_question_regions_enhanced(page_path, known_questions)
        
        # Convert to grayscale and apply preprocessing
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        
        # Apply adaptive thresholding for better shape detection
        thresh = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                                     cv2.THRESH_BINARY_INV, 11, 2)
        
        # Apply morphological operations to connect nearby components
        kernel = np.ones((3, 3), np.uint8)
        thresh = cv2.morphologyEx(thresh, cv2.MORPH_CLOSE, kernel)
        
        # Find contours
        contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        
        for cnt in contours:
            x, y, w, h = cv2.boundingRect(cnt)
            area = cv2.contourArea(cnt)
            
            # Check if this is a valid diagram
            if is_valid_diagram(cnt, w, h, area):
                # Find associated question
                question_num = find_associated_question(y, x, h, question_regions, page_width)
                
                if question_num:
                    # Determine padding based on shape
                    peri = cv2.arcLength(cnt, True)
                    approx = cv2.approxPolyDP(cnt, 0.02 * peri, True)
                    aspect_ratio = float(w) / h
                    
                    # Check if it's circular/oval (many vertices) or rectangular
                    if len(approx) >= 6 and 0.8 < aspect_ratio < 1.2:
                        pad = 65  # Circular/oval shape
                    else:
                        pad = 20    # Rectangular shape
                    
                    # Crop with padding
                    x1 = max(x - pad, 0)
                    y1 = max(y - pad, 0)
                    x2 = min(x + w + pad, img.shape[1])
                    y2 = min(y + h + pad, img.shape[0])
                    
                    cropped = img[y1:y2, x1:x2]
                    
                    # Generate unique filename for multiple diagrams per question
                    existing_count = len(results[question_num])
                    suffix = f"_{chr(97 + existing_count)}" if existing_count > 0 else ""
                    image_filename = f"q{question_num}_diagram{suffix}.png"
                    image_path = os.path.join(output_dir, image_filename)
                    
                    cv2.imwrite(image_path, cropped)
                    results[question_num].append(image_path)
    
    return results


def extract_images_from_document(doc_path):
    """Extract images from the Word document and associate them with questions"""
    # Create temporary directory
    tmpdir = os.path.join(os.path.dirname(doc_path), "temp_processing")
    os.makedirs(tmpdir, exist_ok=True)
    
    try:
        # Convert DOCX to PDF
        pdf_path = convert_docx_to_pdf(doc_path, tmpdir)
        
        # Convert PDF to images
        pages_dir = os.path.join(tmpdir, "pages")
        page_images = convert_pdf_to_images(pdf_path, pages_dir)
        
        # Extract diagrams and associate with questions
        diagrams_dir = os.path.join(os.path.dirname(doc_path), "extracted_diagrams")
        question_diagram_map = extract_diagrams_from_pages(page_images, diagrams_dir)
        
        # Clean up temporary files
        import shutil
        if os.path.exists(tmpdir):
            shutil.rmtree(tmpdir)
        
        return question_diagram_map
        
    except Exception as e:
        print(f"Error extracting images: {e}")
        # Clean up on error
        import shutil
        if os.path.exists(tmpdir):
            shutil.rmtree(tmpdir)
        return {}


def create_arrangement_slide(prs, arrangement_text, direction_text=None, num=0, diagram_path=None):
    """Create a slide with the arrangement text and optional diagram"""
    slide_layout = prs.slide_layouts[5]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    add_logo(slide, prs.slide_width, prs.slide_height)
    
    # Set black background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    # Calculate positioning
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # If there's a diagram, adjust layout
    if diagram_path and os.path.exists(diagram_path):
        # Add the diagram in the left 40% area
        try:
            left = slide_width * 0.42
            top = slide_height * 0.2
            width = slide_width * 0.3
            
            slide.shapes.add_picture(diagram_path, left, top, width=width)
            
            # Adjust text position to right side
            text_left = slide_width * 0.4
            text_width = slide_width * 0.58
        except Exception as e:
            print(f"Warning: Could not add diagram: {e}")
            # Fall back to centered text
            text_left = slide_width * 0.4
            text_width = slide_width * 0.58
    else:
        # Center text if no diagram
        text_left = slide_width * 0.4
        text_width = slide_width * 0.58
    
    text_top = slide_height * 0.050 # Center vertically
    text_height = slide_height * 0.3
    
    # Add text box
    textbox = slide.shapes.add_textbox(
        left=text_left,
        top=text_top,
        width=text_width,
        height=text_height
    )
    
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    # Remove title placeholder if it exists
    for shape in slide.shapes:
        if shape == slide.shapes.title:
            sp = shape
            slide.shapes._spTree.remove(sp._element)
            break

    # Add direction if exists
    if direction_text and num == 0:
        p = text_frame.add_paragraph()
        p.text = direction_text
        p.font.name = 'Arial'
        p.font.size = Pt(25)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)
    
    # Add arrangement text
    p = text_frame.add_paragraph()
    p.text = arrangement_text.strip()
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    add_yellow_border(slide)
    return slide


def create_question_slide(prs, question_data):
    """Create a slide with the MCQ question"""
    slide_layout = prs.slide_layouts[5]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    add_logo(slide, prs.slide_width, prs.slide_height)
    
    # Set black background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    # Calculate positioning
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    text_left = slide_width * 0.4
    text_width = slide_width * 0.58
    text_top = Inches(0.25)
    text_height = slide_height - Inches(1.5)
    
    # Add text box
    textbox = slide.shapes.add_textbox(
        left=text_left,
        top=text_top,
        width=text_width,
        height=text_height
    )
    
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    # Add question number and text
    p = text_frame.add_paragraph()
    p.text = question_data['content']
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.JUSTIFY
    p.space_after = Pt(8)
    
    # Remove title placeholder if it exists
    for shape in slide.shapes:
        if shape == slide.shapes.title:
            sp = shape
            slide.shapes._spTree.remove(sp._element)
            break
    
    # Add options
    for option in question_data['options']:
        p = text_frame.add_paragraph()
        p.text = option
        p.font.name = 'Arial'
        p.font.size = Pt(25)
        p.font.color.rgb = RGBColor(255, 255, 103)  # Yellow for options
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)

    add_yellow_border(slide)
    return slide


def add_yellow_border(slide):
    """Add a yellow border effect to the slide - only top and bottom, 3/4 width"""
    slide_width = Inches(13.33)
    slide_height = Inches(7.5)
    border_width = Inches(0.15)
    
    # Calculate 3/4 of the slide width
    border_length = slide_width * 0.75
    
    start_x = slide_width - border_length
    # Top border (3/4 width from left)
    top_border = slide.shapes.add_shape(
        1,  # Rectangle shape
        start_x,  # Start from right edge
        0,  # Top of slide
        border_length,  # 3/4 of slide width
        border_width    # Border thickness
    )
    top_border.fill.solid()
    top_border.fill.fore_color.rgb = RGBColor(255, 255, 103)
    top_border.line.fill.background()
    
    # Bottom border (3/4 width from left)
    bottom_border = slide.shapes.add_shape(
        1,  # Rectangle shape
        0,  # Start from left edge
        slide_height - border_width,  # Bottom position
        border_length,  # 3/4 of slide width
        border_width    # Border thickness
    )
    bottom_border.fill.solid()
    bottom_border.fill.fore_color.rgb = RGBColor(255, 255, 103)
    bottom_border.line.fill.background()


def convert_word_to_ppt(word_path, ppt_path):
    """Main function to convert Word document to PowerPoint"""
    # Parse Word document
    print("Parsing Word document...")
    questions = parse_word_document(word_path)
    
    # Extract images from document
    print("Extracting diagrams and images...")
    extracted_images = extract_images_from_document(word_path)
    
    # Create PowerPoint presentation
    print("Creating PowerPoint presentation...")
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Track the current direction for all slides
    current_direction = None
    
    # Create slides for each question
    slide_count = 0
    for i, question in enumerate(questions):
        # Extract question number from the content
        question_num_match = re.search(r'^(\*\*)?(\d+)\.', question['content'])
        question_num = int(question_num_match.group(2)) if question_num_match else None
        
        # Update direction if it changes
        if question['direction'] and question['direction'] != current_direction:
            current_direction = question['direction']
        
        # Check if there's a diagram for this question's arrangement
        arrangement_diagram = None
        if question_num and question_num in extracted_images:
            # Use the first diagram for the arrangement slide
            if extracted_images[question_num]:
                arrangement_diagram = extracted_images[question_num][0]
        
        # If there's an arrangement, create a separate slide for it
        if question.get('arrangement'):
            print(f"Creating arrangement slide {slide_count + 1}...")
            create_arrangement_slide(prs, question['arrangement'], current_direction, i, arrangement_diagram)
            slide_count += 1
        
        # Create the question slide
        print(f"Creating question slide {slide_count + 1}...")
        create_question_slide(prs, question)
        slide_count += 1
    
    # Save presentation
    print(f"Saving presentation to {ppt_path}...")
    prs.save(ppt_path)
    print(f"Conversion complete! Created {slide_count} slides.")
    
    # Cleanup extracted images if needed
    extracted_dir = os.path.join(os.path.dirname(word_path), "extracted_diagrams")
    if os.path.exists(extracted_dir):
        import shutil
        shutil.rmtree(extracted_dir, ignore_errors=True)


# Main execution
if __name__ == "__main__":
    # Input and output file paths
    input_file = "mcq2.docx"
    output_file = os.path.join("output2.pptx")
    
    # Check if input file exists
    convert_word_to_ppt(input_file, output_file)