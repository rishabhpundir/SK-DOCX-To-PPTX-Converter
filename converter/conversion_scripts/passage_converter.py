#!/usr/bin/env python3
"""
Word to PowerPoint Converter
Converts MS Word documents containing passages and questions to formatted PowerPoint presentations.
"""

import re
import os
import argparse
from docx import Document
from pptx import Presentation
from django.conf import settings
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


class WordToPowerPointConverter:
    """Main converter class that handles the complete conversion process"""
    
    def __init__(self):
        self.chars_per_slide = 725  # Adjust based on your content needs
        
    def add_logo(self, slide):
        # Add logo to top-left corner
        logo_path = os.path.join(settings.BASE_DIR, "docs", "passage_logo.png")
        logo_left = Inches(-0.4)
        logo_top = Inches(-0.4)
        logo_width = Inches(2.5)
        slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width)
        
    def read_docx_file(self, file_path):
        """Read content from Word document"""
        try:
            doc = Document(file_path)
            content = ""
            for paragraph in doc.paragraphs:
                content += paragraph.text + "\n"
            return content
        except Exception as e:
            print(f"❌ Error reading Word file: {e}")
            return None

    def parse_document_content(self, content):
        """Parse the document content into structured sections"""
        sections = {
            'directions': '',
            'passages': [],
        }
        
        # Extract directions
        directions_match = re.search(r'DIRECTIONS FOR QUESTION.*?(?=PASSAGE)', content, re.DOTALL)
        if directions_match:
            sections['directions'] = directions_match.group().strip()
        
        # Extract passages - handle both -- and – patterns
        passage_pattern = r'(PASSAGE\s*[–-]+\s*[IVX]+)\s*\n(.*?)(?=\nPASSAGE\s*[–-]+\s*[IVX]+|$)'
        passage_matches = re.findall(passage_pattern, content, re.DOTALL)

        for passage_num, passage_text in passage_matches:
            # Clean up passage text
            passage_text = re.sub(r'(\[Extracted.*?\])', r'\1\n\n', passage_text, flags=re.DOTALL)
            passage_content, questions = passage_text.strip().split('\n\n\n\n', 1)
            sections['passages'].append({
                'number': passage_num,
                'content': passage_content,
            })
        
            # Extract questions with their options
            pattern = r'(?=(?:^|\n)\d{1,2}\.\t)'
            question_matches = re.split(pattern, questions.strip())
            question_matches = [part.strip() for part in question_matches if part.strip()]
            sections['passages'][-1]['questions'] = question_matches
        return sections

    def create_title_slide(self, prs, title, subtitle=""):
        """Create a formatted title slide"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        self.add_logo(slide)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        if subtitle:
            subtitle_shape.text = subtitle
        
        # Apply formatting
        self.apply_slide_formatting(slide)
        return slide

    def create_content_slide(self, num, prs, directions, title, content, is_passage=False, is_last_passage_slide=False):
        """Create a formatted content slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        self.add_logo(slide)
        
        if title.strip() != "":
            # Use and modify the default title shape
            title_shape = slide.shapes.title
            
            # Reposition and resize the title shape to match your alignment
            title_shape.left = Inches(5.0)  # Same as content
            title_shape.top = Inches(0.5)   # Start at top
            title_shape.width = Inches(8)   # Same width as content
            title_shape.height = Inches(2)  # Height just for title
            
            # Configure title text formatting
            title_text_frame = title_shape.text_frame
            title_text_frame.word_wrap = True
            
            # Clear the text frame
            title_text_frame.clear()
            
            # Add directions paragraph (font size 16) only if num == 0
            if num == 0 and directions.strip():
                directions_paragraph = title_text_frame.paragraphs[0]
                directions_paragraph.text = directions
                directions_paragraph.alignment = PP_ALIGN.JUSTIFY
                directions_paragraph.runs[0].font.size = Pt(20)
                directions_paragraph.runs[0].font.color.rgb = RGBColor(255, 255, 255)
                
                
                # Add title paragraph (font size 25)
                title_paragraph = title_text_frame.add_paragraph()
            else:
                # If num != 0 or no directions, use the first paragraph for title
                title_paragraph = title_text_frame.paragraphs[0]
            
            # Set title text and formatting
            title_paragraph.text = title
            title_paragraph.alignment = PP_ALIGN.JUSTIFY
            title_paragraph.runs[0].font.size = Pt(22)
            title_paragraph.runs[0].font.color.rgb = RGBColor(255, 255, 255)
            
        else:
            # Remove the default title shape if no title
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    sp = shape
                    slide.shapes._spTree.remove(sp._element)
                    break
        
        if is_passage:
            # For passage slides, use custom layout on the right half
            self.create_passage_content(directions, slide, title, content)
            # Add "Continued" footer only on continuation slides, not on last slide
            if not is_last_passage_slide:
                left = Inches(11.25)  # Align with passage content
                top = Inches(6.80)
                width = Inches(8)
                height = Inches(0.5)
                footer = slide.shapes.add_textbox(left, top, width, height)
                text_frame = footer.text_frame
                text_frame.text = "Continued"
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.RIGHT
                p.font.color.rgb = RGBColor(255, 255, 255)
                run = p.runs[0]
                run.font.size = Pt(18)
                run.font.name = "Arial"
        else:
            # For other content, use standard layout
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = content
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.RIGHT
        
        # Apply formatting
        self.apply_slide_formatting(slide)
        return slide

    def create_passage_content(self, directions, slide, title, content):
        """Create passage content in the right half of the slide"""
        # Remove default content placeholder if it exists
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape._sp
            sp.getparent().remove(sp)
        
        # Add new text box on the right half of the slide
        left = Inches(5.0)  # Start from middle of slide
        if title.strip() != "":
                top = Inches(2.25)
        else:
                top = Inches(0.5)
        width = Inches(8)  # Right half width
        height = Inches(5.5)  # Increased from 5.0 to use more vertical space
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        # Configure text frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        # Add content
        p = text_frame.paragraphs[0]
        p.text = content
        p.alignment = PP_ALIGN.JUSTIFY
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
        p.font.name = 'Arial'
        
    def split_mcq_list(self, mcq_list):
        """
        Split MCQ strings that use either tabs or newlines as separators.
        Supports numbers, letters, and Roman numerals with various formats.
        """
        result = []
        
        for item in mcq_list:
            # Check if item contains tabs or newlines (needs splitting)
            if '\t' in item or '\n' in item:
                # Replace tabs after various numbering patterns with space
                # Patterns supported:
                # - Numbers: 1.\t, 1)\t, (1)\t
                # - Letters: A)\t, (A)\t
                # - Roman numerals: I.\t, II.\t, IV.\t, (I)\t, I)\t, etc.
                roman_pattern = 'M{0,3}(CM|CD|D?C{0,3})(XC|XL|L?X{0,3})(IX|IV|V?I{0,3})'
                pattern = rf'(\d+\.|\d+\)|\(\d+\)|[A-Z]\)|\([A-Z]\)|{roman_pattern}\.|\({roman_pattern}\)|{roman_pattern}\))\t+'
                item = re.sub(pattern, r'\1 ', item)
                
                # Split by both newlines and remaining tabs
                parts = re.split(r'[\n\t]+', item)
                
                # Clean and filter parts
                parts = [part.strip() for part in parts if part.strip()]
                result.extend(parts)
            else:
                # Item doesn't need splitting
                result.append(item.strip())
        
        return result

    def create_questions_slide(self, prs, questions_list):
        """Create a formatted slide with questions"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        self.add_logo(slide)
        
        # Remove title shape
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                sp = shape
                slide.shapes._spTree.remove(sp._element)
                break
        
        # Remove default content placeholder
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape._sp
            sp.getparent().remove(sp)
        
        # Add new text box starting 5 inches from left
        left = Inches(5.0)
        top = Inches(0.75)  # Start from top
        width = Inches(8)
        height = Inches(6.25)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        
        questions_list = self.split_mcq_list(questions_list)
        for i, question in enumerate(questions_list):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = question
            p.font.size = Pt(21)
            p.space_after = Pt(10)
            p.alignment = PP_ALIGN.LEFT  # Left align within the right-positioned textbox
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = 'Arial'
        
        # Apply formatting
        self.apply_slide_formatting(slide)
        return slide

    def split_passage_content(self, content):
        """Split passage content into multiple chunks if needed"""
        # Split content into sentences to avoid breaking mid-sentence
        sentences = content.split('.')
        sentences = [s + '.'if s[-1] not in ('.', ']') else s for s in [s for s in sentences if s.strip()] ]
        
        # Group sentences into slide chunks
        slide_contents = []
        current_chunk = []
        current_length = 0
        
        for sentence in sentences:
            self.chars_per_slide = 725
            if len(slide_contents) == 0:
                self.chars_per_slide = 550
            if current_length + len(sentence) > self.chars_per_slide and current_chunk:
                slide_contents.append(''.join(current_chunk))
                current_chunk = [sentence]
                current_length = len(sentence)
            else:
                current_chunk.append(sentence)
                current_length += len(sentence)
        
        if current_chunk:
            slide_contents.append(''.join(current_chunk))
        return slide_contents

    def create_passage_slides(self, prs, passage, directions):
        """Create one or more slides for a passage with proper content distribution"""
        passage_num = passage['number']
        content = passage['content']
        
        # Split content if it's too long
        slide_contents = self.split_passage_content(content)
        
        created_slides = []
        for i, slide_content in enumerate(slide_contents):
            if i == 0:    
                title = passage_num
            else:
                title = ""
            
            # Check if this is the last slide for this passage
            is_last_slide = (i == len(slide_contents) - 1)
            slide_content = slide_content.strip()
            
            slide = self.create_content_slide(i, prs, directions, title, slide_content, is_passage=True, is_last_passage_slide=is_last_slide)
            created_slides.append(slide)
        
        return created_slides

    def apply_slide_formatting(self, slide):
        """Apply comprehensive formatting to a slide"""
        # Set black background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black
        
        # Format title
        if slide.shapes.title:
            title_frame = slide.shapes.title.text_frame
            for paragraph in title_frame.paragraphs:
                # paragraph.font.size = Pt(22)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
                paragraph.font.bold = True
                paragraph.font.name = 'Arial'
                paragraph.alignment = PP_ALIGN.JUSTIFY
        
        # Format content shapes
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                text_frame = shape.text_frame
                
                # Check if this is an MCQ slide
                is_mcq = False
                if slide.shapes.title and slide.shapes.title.text:
                    if "QUESTION" in slide.shapes.title.text.upper():
                        is_mcq = True
                
                for paragraph in text_frame.paragraphs:
                    # Set font size based on content type
                    if not is_mcq and slide.shapes.title and "PASSAGE" in slide.shapes.title.text.upper():
                        paragraph.font.size = Pt(22)
                        paragraph.alignment = PP_ALIGN.JUSTIFY
                    else:
                        paragraph.font.size = Pt(22)
                        paragraph.alignment = PP_ALIGN.LEFT
                    
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
                    paragraph.font.name = 'Arial'
        
        # Add red border
        self.add_red_border(slide)

    def add_red_border(self, slide):
        """Add a red border effect to the slide - only top and bottom, 3/4 width"""
        slide_width = Inches(13.33)
        slide_height = Inches(7.5)
        border_width = Inches(0.1)
        
        # Calculate 3/4 of the slide width
        border_length = slide_width * 0.75
        
        start_x = slide_width - border_length
        # Top border (3/4 width from left)
        top_border = slide.shapes.add_shape(
            1,  # Rectangle shape
            start_x,  # Start from right edge
            0,  # Top of slide
            border_length,  # 3/4 of slide width
            border_width    # Border thickness
        )
        top_border.fill.solid()
        top_border.fill.fore_color.rgb = RGBColor(255, 0, 0)
        top_border.line.fill.background()
        
        # Bottom border (3/4 width from left)
        bottom_border = slide.shapes.add_shape(
            1,  # Rectangle shape
            0,  # Start from left edge
            slide_height - border_width,  # Bottom position
            border_length,  # 3/4 of slide width
            border_width    # Border thickness
        )
        bottom_border.fill.solid()
        bottom_border.fill.fore_color.rgb = RGBColor(255, 0, 0)
        bottom_border.line.fill.background()

    def convert(self, input_docx_path, output_pptx_path=None):
        """Main conversion method"""
        # Set default output path if not provided
        if output_pptx_path is None:
            base_name = os.path.splitext(os.path.basename(input_docx_path))[0]
            output_pptx_path = f"{base_name}_formatted.pptx"
        
        # Validate input file
        if not os.path.exists(input_docx_path):
            print(f"❌ Error: Input file '{input_docx_path}' not found!")
            return False
        
        # Read Word document
        print(f"📖 Reading Word document: {input_docx_path}")
        content = self.read_docx_file(input_docx_path)
        
        if content is None:
            print("❌ Failed to read Word document.")
            return False
        
        # Create presentation with custom dimensions
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Parse content
        print("🔍 Parsing document content...")
        sections = self.parse_document_content(content)
        
        # Create title slide
        self.create_title_slide(prs, "Section I - English", "Reading Comprehension Test")
        
        # Create slides for each passage
        print(f"📄 Processing {len(sections['passages'])} passages...")
        for passage in sections['passages']:
            self.create_passage_slides(prs, passage, sections['directions'])
        
            # Create question slides
            print(f"❓ Processing {len(passage['questions'])} questions...")
            questions_per_slide = 1
            question_groups = [passage['questions'][i:i + questions_per_slide] 
                            for i in range(0, len(passage['questions']), questions_per_slide)]
            
            for i, question_group in enumerate(question_groups):
                start_q = i * questions_per_slide + 1
                self.create_questions_slide(prs, question_group)
        
        # Save presentation
        try:
            prs.save(output_pptx_path)
            print(f"✅ Formatted PowerPoint presentation saved as: {output_pptx_path}")
            print(f"📊 Created {len(prs.slides)} slides total")
            return True
        except Exception as e:
            print(f"❌ Error saving presentation: {e}")
            return False


def main():
    """Main function with command-line argument parsing"""
    parser = argparse.ArgumentParser(
        description="Convert MS Word documents to formatted PowerPoint presentations",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python word_to_ppt_converter.py input.docx
  python word_to_ppt_converter.py input.docx -o output.pptx
  python word_to_ppt_converter.py /path/to/document.docx -o /path/to/output.pptx
        """
    )
    
    parser.add_argument(
        'input_file',
        help='Path to the input Word document (.docx)'
    )
    
    parser.add_argument(
        '-o', '--output',
        dest='output_file',
        help='Path for the output PowerPoint file (.pptx). If not specified, will be generated automatically.'
    )
    
    parser.add_argument(
        '--chars-per-slide',
        type=int,
        default=800,
        help='Maximum characters per passage slide (default: 800)'
    )
    
    args = parser.parse_args()
    
    # Validate input file extension
    if not args.input_file.lower().endswith('.docx'):
        print("❌ Error: Input file must be a Word document (.docx)")
        return False
    
    # Create converter instance
    converter = WordToPowerPointConverter()
    converter.chars_per_slide = args.chars_per_slide
    
    # Perform conversion
    print("🚀 Starting Word to PowerPoint conversion...")
    success = converter.convert(args.input_file, args.output_file)
    
    if success:
        print("🎉 Conversion completed successfully!")
        return True
    else:
        print("❌ Conversion failed!")
        return False


if __name__ == "__main__":
    main()