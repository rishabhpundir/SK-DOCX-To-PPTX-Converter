import re
import os
from docx import Document
from pptx import Presentation
from django.conf import settings
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_logo(slide, slide_width, slide_height):
    # Add logo to top-left corner
    logo_path = os.path.join(settings.BASE_DIR, "docs", "mcq_logo.png")
    bg_logo_path = os.path.join(settings.BASE_DIR, "docs", "bg_logo.png")
    logo_left = Inches(0.2)
    logo_top = Inches(0.2)
    logo_width = Inches(1.0)  # Adjust as needed
    slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width)

    img_width = Inches(3.5)
    img_height = Inches(3)

    # Center position
    left = (slide_width - img_width) / 2
    top = (slide_height - img_height) / 2

    # Add image
    slide.shapes.add_picture(bg_logo_path, left, top, width=img_width, height=img_height)
    

def add_yellow_border(slide):
    """Add a yellow border effect to the slide - only top and bottom, 3/4 width"""
    slide_width = Inches(13.33)
    slide_height = Inches(7.5)
    border_width = Inches(0.15)
    
    # Calculate 3/4 of the slide width
    border_length = slide_width * 0.75
    
    start_x = slide_width - border_length
    # Top border (3/4 width from left)
    top_border = slide.shapes.add_shape(
        1,  # Rectangle shape
        start_x,  # Start from right edge
        0,  # Top of slide
        border_length,  # 3/4 of slide width
        border_width    # Border thickness
    )
    top_border.fill.solid()
    top_border.fill.fore_color.rgb = RGBColor(255, 255, 103)
    top_border.line.fill.background()
    
    # Bottom border (3/4 width from left)
    bottom_border = slide.shapes.add_shape(
        1,  # Rectangle shape
        0,  # Start from left edge
        slide_height - border_width,  # Bottom position
        border_length,  # 3/4 of slide width
        border_width    # Border thickness
    )
    bottom_border.fill.solid()
    bottom_border.fill.fore_color.rgb = RGBColor(255, 255, 103)
    bottom_border.line.fill.background()
    

def parse_word_document(doc_path):
    """Parse the Word document and extract questions with their directions"""
    doc = Document(doc_path)
    
    content_blocks = []
    current_direction = None
    current_question = None
    
    text_info = False
    for para in doc.paragraphs:
        text = para.text.strip()

        if not text:
            continue
        
        if text_info:
            text_info = False
            content_blocks.append({
                    'type': 'info',
                    'direction': current_direction,
                    'content': text
                })
            
            
        # Check if it's a direction
        if "Directions for questions" in text or "DIRECTIONS:" in text:
            text_info = True
            # Save previous question if exists
            if current_question:
                content_blocks.append({
                    'type': 'question',
                    'direction': current_direction,
                    'content': current_question
                })
                current_question = None
            
            # Start new direction
            current_direction = text
            # Get the direction details from next paragraph
            
        # Check if it's a question number
        elif re.match(r'^\d{1,2}+\.', text):
            # Save previous question if exists
            if current_question:
                content_blocks.append({
                    'type': 'question',
                    'direction': current_direction,
                    'content': current_question
                })
            
            # Start new question
            current_question = text
            
        # If we're in a question, append the options
        elif current_question:
            current_question += '\n' + text
    
    # Don't forget the last question
    if current_question:
        content_blocks.append({
            'type': 'question',
            'direction': current_direction,
            'content': current_question
        })
    return content_blocks


def split_mcq_list(mcq_list):
    """Split MCQ strings while keeping option numbers with their values."""
    result = []
    for item in mcq_list:
        if '\t' in item:
            # Replace tabs between option parts with space, but keep tabs between different options
            # This regex finds patterns like (1)\t and replaces the tab with space
            item = re.sub(r'(\d+\.|\d+\)|\(\d+\))\t+', r'\1 ', item)
            # Now split by remaining tabs (which separate different options)
            parts = [part.strip() for part in re.split(r'\t+', item) if part.strip()]
            result.extend(parts)
        else:
            result.append(item.strip())
    return result


def create_slide(prs, question_data):
    """Create a slide with the MCQ question"""
    slide_layout = prs.slide_layouts[5]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    add_logo(slide, prs.slide_width, prs.slide_height)

    for shape in slide.shapes:
        if shape == slide.shapes.title:
            sp = shape
            slide.shapes._spTree.remove(sp._element)
            break
    
    # Set black background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    # Calculate positioning (60% of slide width, starting from 40%)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    text_left = slide_width * 0.4  # Start at 40% from left
    text_width = slide_width * 0.58  # Use 58% width (leaving 2% margin)
    text_top = Inches(1)  # Start 1 inch from top
    text_height = slide_height - Inches(1.5)  # Leave some bottom margin
    
    # Add text box
    textbox = slide.shapes.add_textbox(
        left=text_left,
        top=text_top,
        width=text_width,
        height=text_height
    )
    
    text_frame = textbox.text_frame
    text_frame.clear()  # Clear default text
    text_frame.word_wrap = True
    
    # Add direction if exists
    if question_data['direction'] and question_data['type'] == 'info':
        p = text_frame.add_paragraph()
        p.text = question_data['direction']
        p.font.name = 'Arial'
        p.font.size = Pt(25)
        p.font.color.rgb = RGBColor(255, 255, 255)  # White
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        # Add empty line after direction
        p = text_frame.add_paragraph()
        p.text = ""
    
    # Process question content
    if question_data['type'] == 'question':
        lines_list = question_data['content'].split('\n')
        lines = split_mcq_list(lines_list)
    else:
        lines = [question_data['content']]
    
    for i, line in enumerate(lines):
        if i == 0:  # First paragraph (no need to add new)
            if question_data['direction']:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line
        p.font.name = 'Arial'
        p.font.size = Pt(25)
        
        # Check if this is a question line (starts with number) or contains "?" 
        if re.match(r'^\d{1,2}+\.', line) or '?' in line or question_data['type'] == 'info':
            p.font.color.rgb = RGBColor(255, 255, 255)  # White for questions
            p.alignment = PP_ALIGN.JUSTIFY
        else:
            # This is likely an option line
            p.font.color.rgb = RGBColor(255, 255, 103)  # Yellow for options
            p.alignment = PP_ALIGN.LEFT
        
        
        # Add some spacing between lines
        p.space_after = Pt(6)
        
    add_yellow_border(slide)
    return slide

def convert_word_to_ppt(word_path, ppt_path):
    """Main function to convert Word document to PowerPoint"""
    # Parse Word document
    print("Parsing Word document...")
    questions = parse_word_document(word_path)
    
    # Create PowerPoint presentation
    print("Creating PowerPoint presentation...")
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Create slides for each question
    for i, question in enumerate(questions):
        print(f"Creating slide {i+1}...")
        create_slide(prs, question)
    
    # Save presentation
    print(f"Saving presentation to {ppt_path}...")
    prs.save(ppt_path)
    print("Conversion complete!")

# Main execution
if __name__ == "__main__":
    # Input and output file paths
    input_file = "mcq3.docx"
    output_file = "output3.pptx"
    
    # Check if input file exists
    if not os.path.exists(input_file):
        print(f"Error: Input file '{input_file}' not found!")
    else:
        convert_word_to_ppt(input_file, output_file)
        print(f"Successfully created {output_file}")