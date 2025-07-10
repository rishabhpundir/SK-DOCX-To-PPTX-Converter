#!/usr/bin/env python3
"""
MCQ Word to PowerPoint Converter
Converts Word documents containing MCQs directly to formatted PowerPoint presentations
"""

import re
import os
import cv2
import uuid
import shutil
import subprocess
import pytesseract
from PIL import Image
from pathlib import Path
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
from pptx.util import Inches, Pt
from django.conf import settings
from pptx.enum.text import PP_ALIGN
from collections import defaultdict
from pptx.dml.color import RGBColor
from pdf2image import convert_from_path

class MCQConverter:
    def __init__(self):
        # MCQ parsing patterns
        self.mcq_pattern = r'\*\*(\d+)\.\*\*\s*(.*?)(?=\*\*\d+\.\*\*|\Z)'
        self.option_pattern = r'\\?\((\d+)\)\s*([^\\(]+?)(?=\\?\(\d+\)|$)'
        self.SUPERSCRIPTS = {
            '0': '⁰', '1': '¹', '2': '²', '3': '³',
            '4': '⁴', '5': '⁵', '6': '⁶',
            '7': '⁷', '8': '⁸', '9': '⁹'
        }
        
        # Color definitions
        self.bg_color = RGBColor(0, 0, 0)  # Black background
        self.question_color = RGBColor(255, 255, 255)  # White for questions
        self.option_color = RGBColor(255, 255, 103)  # Gold/Yellow for options
        
        # Slide dimensions (16:9)
        self.slide_width = Inches(13.33)
        self.slide_height = Inches(7.5)
        
        # Layout settings
        self.left_margin = Inches(5.33)  # 40% of slide width
        self.content_width = Inches(7.5)  # 60% of slide width
        self.top_margin = Inches(0.25)
        self.content_height = Inches(6.5)
        
    def add_logo(self, slide, slide_width, slide_height):
        # Add logo to top-left corner
        logo_path = os.path.join(settings.BASE_DIR, "docs", "mcq_logo.png")
        bg_logo_path = os.path.join(settings.BASE_DIR, "docs", "bg_logo.png")
        logo_left = Inches(0.2)
        logo_top = Inches(0.2)
        logo_width = Inches(1.0)  # Adjust as needed
        slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width)

        img_width = Inches(3.5)
        img_height = Inches(3)

        # Center position
        left = (slide_width - img_width) / 2
        top = (slide_height - img_height) / 2

        # Add image
        slide.shapes.add_picture(bg_logo_path, left, top, width=img_width, height=img_height)

        
    def add_yellow_border(self, slide):
        """Add a yellow border effect to the slide - only top and bottom, 3/4 width"""
        slide_width = Inches(13.33)
        slide_height = Inches(7.5)
        border_width = Inches(0.15)
        
        # Calculate 3/4 of the slide width
        border_length = slide_width * 0.75
        
        start_x = slide_width - border_length
        # Top border (3/4 width from left)
        top_border = slide.shapes.add_shape(
            1,  # Rectangle shape
            start_x,  # Start from right edge
            0,  # Top of slide
            border_length,  # 3/4 of slide width
            border_width    # Border thickness
        )
        top_border.fill.solid()
        top_border.fill.fore_color.rgb = RGBColor(255, 255, 103)
        top_border.line.fill.background()
        
        # Bottom border (3/4 width from left)
        bottom_border = slide.shapes.add_shape(
            1,  # Rectangle shape
            0,  # Start from left edge
            slide_height - border_width,  # Bottom position
            border_length,  # 3/4 of slide width
            border_width    # Border thickness
        )
        bottom_border.fill.solid()
        bottom_border.fill.fore_color.rgb = RGBColor(255, 255, 103)
        bottom_border.line.fill.background()
    
    def extract_text_with_positions(self, docx_path):
        """Extract text and try to identify question numbers"""
        doc = Document(docx_path)
        questions = []
        current_question = None
        
        for para in doc.paragraphs:
            text = para.text.strip()
            # Look for question numbers at the start of paragraphs
            match = re.match(r'^(\d{1,2}+)\.\s*', text)
            if match:
                question_num = int(match.group(1))
                questions.append({
                    'number': question_num,
                    'text': text
                })
                current_question = question_num
        
        return questions

    def convert_docx_to_pdf_python(self, docx_path, output_dir):
        """
        Converts a DOCX file to PDF using LibreOffice (headless).
        Works cross-platform if LibreOffice is installed.
        """
        # Ensure output dir exists
        os.makedirs(output_dir, exist_ok=True)

        # Run the conversion
        try:
            subprocess.run([
                "soffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", output_dir,
                docx_path
            ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        except FileNotFoundError:
            raise EnvironmentError("LibreOffice (`soffice`) not found. Please install it and ensure it's in your PATH.")
        except subprocess.CalledProcessError as e:
            raise RuntimeError(f"LibreOffice failed to convert the file: {e}")

        # Construct output PDF path
        pdf_name = Path(docx_path).stem + ".pdf"
        pdf_path = os.path.join(output_dir, pdf_name)

        # Validate conversion
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"Expected PDF not found at {pdf_path}")

        print(f"[✓] Converted DOCX to PDF using LibreOffice: {pdf_path}")
        return pdf_path

    def convert_pdf_to_images(self, pdf_path, image_dir):
        """Convert PDF pages to images"""
        os.makedirs(image_dir, exist_ok=True)
        images = convert_from_path(pdf_path, dpi=300)
        page_paths = []

        for i, img in enumerate(images):
            img_path = os.path.join(image_dir, f"page_{i+1}.png")
            img.save(img_path, "PNG")
            page_paths.append(img_path)
        print(f"[✓] Rendered {len(page_paths)} pages to {image_dir}")
        return page_paths

    def detect_question_regions(self, page_path):
        """Detect regions containing question numbers using OCR"""
        img = cv2.imread(page_path)
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        
        # Use OCR to detect text and positions
        try:
            # Get OCR data with bounding boxes
            data = pytesseract.image_to_data(gray, output_type=pytesseract.Output.DICT)
            
            question_regions = []
            
            for i in range(len(data['text'])):
                text = str(data['text'][i]).strip()
                # Look for question numbers
                if re.match(r'^\d{1,2}+\.$', text):
                    question_num = int(text[:-1])
                    x, y, w, h = data['left'][i], data['top'][i], data['width'][i], data['height'][i]
                    
                    # Store question number and its position
                    question_regions.append({
                        'number': question_num,
                        'x': x,
                        'y': y,
                        'width': w,
                        'height': h
                    })
                    
            return question_regions
        except Exception as e:
            print(f"[!] OCR error: {e}")
            return []

    def find_closest_question(self, image_y, image_x, question_regions, page_width):
        """Find which question number an image belongs to based on position"""
        if not question_regions:
            return None

        # Find the question that is above and closest to the image
        closest_question = None
        min_distance = float('inf')
        
        for q in question_regions:
            # Only consider questions that are above the image
            if q['y'] < image_y:
                distance = image_y - q['y']
                if distance < min_distance:
                    min_distance = distance
                    closest_question = q['number']
                    
        return closest_question

    def extract_images_with_questions(self, page_paths, output_dir, area_threshold=10000):
        """Extract images and associate them with question numbers"""
        os.makedirs(output_dir, exist_ok=True)
        results = defaultdict(lambda: None)  # Dictionary to store question -> image mapping
        
        for page_idx, page_path in enumerate(page_paths):
            img = cv2.imread(page_path)
            page_height, page_width = img.shape[:2]
            
            # Detect question regions on this page
            question_regions = self.detect_question_regions(page_path)
            
            # Extract images using existing logic
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            blur = cv2.GaussianBlur(gray, (5, 5), 0)
            _, thresh = cv2.threshold(blur, 200, 255, cv2.THRESH_BINARY_INV)
            
            contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            for cnt in contours:
                x, y, w, h = cv2.boundingRect(cnt)
                area = w * h
                
                if area > area_threshold:
                    # Determine padding based on shape
                    peri = cv2.arcLength(cnt, True)
                    approx = cv2.approxPolyDP(cnt, 0.04 * peri, True)
                    aspect_ratio = float(w) / h
                    
                    if len(approx) >= 6 and 0.8 < aspect_ratio < 1.2:
                        pad = 140 # Circular/oval shape
                    else:
                        pad = 15   # Rectangular shape
                    
                    # Find which question this image belongs to
                    image_center_y = y + h // 2
                    image_center_x = x + w // 2
                    question_num = self.find_closest_question(image_center_y, image_center_x, 
                                                       question_regions, page_width)
                    
                    if question_num:
                        # Crop and save the image
                        x1 = max(x - pad, 0)
                        y1 = max(y - pad, 0)
                        x2 = min(x + w + pad, img.shape[1])
                        y2 = min(y + h + pad, img.shape[0])
                        
                        cropped = img[y1:y2, x1:x2]
                        
                        # Save image with question number in filename
                        timestamp = str(uuid.uuid4()).replace('-', '')[:8]
                        image_filename = f"question_{question_num}_image_{timestamp}.png"
                        image_path = os.path.join(output_dir, image_filename)
                        cv2.imwrite(image_path, cropped)
                        
                        # Store in results
                        if results[question_num] is None:
                            results[question_num] = [image_path]
                        else:
                            results[question_num].append(image_path)
                        print(f"[✓] Found image for Question {question_num}")
        return results

    def extract_mcq_images(self, docx_path):
        """Main function to extract images from MCQ document"""
        output_root = os.path.join(settings.BASE_DIR, "media", "extraction")
        
        if os.path.exists(output_root):
            shutil.rmtree(output_root)
            print(f"[✓] Cleared cache...")
            
        tmpdir = os.path.join(output_root, "temp")
        os.makedirs(tmpdir, exist_ok=True)
        
        # Step 1: Extract question numbers from text
        questions = self.extract_text_with_positions(docx_path)
        all_question_numbers = [q['number'] for q in questions]
        
        # Step 2: Convert DOCX to PDF using Python package
        pdf_path = self.convert_docx_to_pdf_python(docx_path, tmpdir)
        assert os.path.exists(pdf_path), f"PDF not found at: {pdf_path}"
        
        # Step 3: Convert PDF to images
        image_dir = os.path.join(output_root, "pages")
        page_images = self.convert_pdf_to_images(pdf_path, image_dir)
        
        # Step 4: Extract images and associate with questions
        images_dir = os.path.join(output_root, "extracted_images")
        question_image_map = self.extract_images_with_questions(page_images, images_dir)
        
        # Step 5: Create final result list
        result = []
        
        # Add all questions, including those without images
        for question_num in sorted(all_question_numbers):
            if question_num in question_image_map:
                result.append({
                    "name": f"{question_num}",
                    "image": question_image_map[question_num]
                })
            else:
                result.append({
                    "name": f"{question_num}",
                    "image": None
                })

        print(f"\n[✓] Extraction complete. Found {len(question_image_map)} images out of {len(all_question_numbers)} questions.")
        return result
    

    def convert_docx_to_html(self, docx_path):
        output_dir = os.path.join("test_files" , "temp")
        os.makedirs(output_dir, exist_ok=True)
        html_file = os.path.join(output_dir, os.path.splitext(os.path.basename(docx_path))[0] + ".html")

        subprocess.run([
            "soffice", "--headless", "--convert-to", "html:XHTML Writer File:UTF8", "--outdir", output_dir, docx_path
        ], check=True)

        if not os.path.exists(html_file):
            raise FileNotFoundError("HTML file not generated")
        return html_file


    def replace_mathml_superscripts(self, soup):
        # Handle <msup> tags as before
        for msup in soup.find_all('msup'):
            base_elem = msup.find(['mi', 'mrow'])
            exp_elem = msup.find('mn')

            base_text = ''.join(base_elem.stripped_strings) if base_elem else ''
            exp_text = ''.join(exp_elem.stripped_strings) if exp_elem else ''

            superscript = ''.join(self.SUPERSCRIPTS.get(ch, ch) for ch in exp_text)
            msup.replace_with(f"{base_text}{superscript}")

        # Handle <span> containing <math> with numeric sibling
        for span in soup.find_all('span'):
            if span.find('math'):
                next_sibling = span.find_next_sibling('span')
                if next_sibling and next_sibling.string and next_sibling.string.strip().isdigit():
                    digits = next_sibling.string.strip()
                    superscript = ''.join(self.SUPERSCRIPTS.get(ch, ch) for ch in digits)

                    math_text = ''.join(span.stripped_strings)
                    next_sibling.decompose()
                    span.replace_with(f"{math_text}{superscript}")

        # NEW: Replace spans where class contains "text-T" and text is 2 or 3
        for span in soup.find_all('span', class_=re.compile(r'text-T')):
            if span.string and span.string.strip() in ['2', '3']:
                span.string.replace_with(self.SUPERSCRIPTS.get(span.string.strip(), span.string.strip()))
        return soup


    def replace_mathml_elements(self, soup):
        # Superscripts: <msup>
        for msup in soup.find_all('msup'):
            base_elem = msup.find(['mi', 'mrow'])
            exp_elem = msup.find('mn')

            base_text = ''.join(base_elem.stripped_strings) if base_elem else ''
            exp_text = ''.join(exp_elem.stripped_strings) if exp_elem else ''

            # Convert digits in exponent to superscript Unicode
            superscript = ''.join(self.SUPERSCRIPTS.get(ch, ch) for ch in exp_text)

            msup.replace_with(f"{base_text}{superscript}")

        # Square roots: <msqrt>
        for msqrt in soup.find_all('msqrt'):
            content = ''.join(msqrt.stripped_strings)
            msqrt.replace_with(f"√{content}")

        # Fractions: <mfrac>
        for mfrac in soup.find_all('mfrac'):
            num_elem = mfrac.find_all(['mn', 'mi'])
            numerator = ''.join(num_elem[0].stripped_strings) if len(num_elem) > 0 else ''
            denominator = ''.join(num_elem[1].stripped_strings) if len(num_elem) > 1 else ''
            mfrac.replace_with(f"({numerator}/{denominator})")
        return soup


    def remove_tables_and_graphics(self, soup):
        # Remove all <table> tags and their contents
        for table in soup.find_all('table'):
            table.decompose()

        # Remove all tags with a class starting with "graphic-"
        for tag in soup.find_all(True):
            if not tag.name:
                continue
            class_list = tag.get('class')
            if class_list and any(cls.startswith('graphic-') for cls in class_list):
                tag.decompose()
        return soup


    def split_mcq_blocks(self, text):
        pattern = re.compile(r'\n{3,}\s*\d{1,2}\.\s*\n{3,}')
        matches = list(pattern.finditer(text))

        # Build list of full MCQ blocks
        blocks = []
        for i, match in enumerate(matches):
            start = match.start()
            end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
            blocks.append(text[start:end].strip())
        return blocks


    def parse_html(self, html_file_path):
        with open(html_file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file, 'html.parser')
        
        exp_soup = self.replace_mathml_superscripts(soup)
        tables_removed_soup = self.remove_tables_and_graphics(exp_soup)
        final_soup = self.replace_mathml_elements(tables_removed_soup)
        full_text = final_soup.get_text(separator='\n\n\n', strip=True).replace("\xa0", "")
        if "\n\n\n." in full_text:
            full_text = full_text.replace("\n\n\n.", ".")
        blocks = self.split_mcq_blocks(full_text)

        mcqs = []
        current_mcq = {}
        for index, block in enumerate(blocks, start=1):
            lines = block.split('\n\n\n')
            combined = ' '.join(line.strip() for line in lines if line.strip())
            option_parts = re.split(r'(?=\(\d+\))', combined)
            question_text = option_parts[0].strip()
            options = [opt.strip() for opt in option_parts[1:] if opt.strip()]

            # Build the dictionary
            current_mcq = {
                "number": index,
                "question": question_text,
                "options": options
            }
            mcqs.append(current_mcq)
        return mcqs
    
    
    def set_slide_background(self, slide):
        """Set slide background to black"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color
    
    def add_question_directive(self, slide):
        """Add the directive text at the top of the slide"""
        directive_box = slide.shapes.add_textbox(
            self.left_margin,
            Inches(0.2),
            self.content_width,
            Inches(0.8)
        )
        
        text_frame = directive_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = "DIRECTIONS: Select the correct alternative from the given choices."
        p.alignment = PP_ALIGN.LEFT
        p.font.name = 'Arial'
        p.font.size = Pt(20)
        p.font.color.rgb = self.question_color
        p.font.bold = True
    
    def create_formatted_slide(self, prs, mcq, mcq_image, is_first_slide=False):
        """Create a single formatted slide with MCQ content and image"""
        # Create new slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        self.add_logo(slide, prs.slide_width, prs.slide_height)
        
        # Set background
        self.set_slide_background(slide)
        
        # Add directive to first slide
        if is_first_slide:
            self.add_question_directive(slide)
            top_margin = Inches(1.0)
        else:
            top_margin = self.top_margin
        
        # Calculate available space for content
        available_height = self.slide_height - top_margin - Inches(0.2)
        
        # Create formatted text box
        text_box = slide.shapes.add_textbox(
            self.left_margin,
            top_margin,
            self.content_width,
            available_height
        )
        
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        
        # Add question text
        p = text_frame.paragraphs[0]
        p.text = f"{mcq['question'].replace("\t", " ")}"
        p.alignment = PP_ALIGN.JUSTIFY_LOW
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.color.rgb = self.question_color
        p.font.bold = False
        p.space_after = Pt(12)
        
        # Add options
        for option in mcq['options']:
            p = text_frame.add_paragraph()
            p.text = option
            p.alignment = PP_ALIGN.LEFT
            p.font.name = 'Arial'
            p.font.size = Pt(18)
            p.font.color.rgb = self.option_color
            p.font.bold = False
            p.space_after = Pt(6)
        
        # Add image if available
        if mcq_image and mcq_image.get('image'):
            image_paths = mcq_image['image']
            for image_path in image_paths:
                if os.path.exists(image_path):
                    try:
                        # Open image to get dimensions
                        with Image.open(image_path) as img:
                            img_width, img_height = img.size
                            aspect_ratio = img_width / img_height
                        
                        # Image should take up bottom 1/3 of text box (2.25 inches)
                        max_image_height = Inches(2.25)
                        max_image_width = self.content_width - Inches(0.4)  # Leave some margin
                        
                        # Scale image to fit while maintaining aspect ratio
                        if (max_image_width / max_image_height) > aspect_ratio:
                            # Height is the limiting factor
                            image_height = max_image_height
                            image_width = image_height * aspect_ratio
                        else:
                            # Width is the limiting factor
                            image_width = max_image_width
                            image_height = image_width / aspect_ratio
                        
                        # Ensure image doesn't exceed maximum dimensions
                        if image_height > max_image_height:
                            image_height = max_image_height
                            image_width = image_height * aspect_ratio
                        
                        # Calculate vertical position for image
                        # Place it in the bottom 1/3 of the text box
                        text_box_bottom = top_margin + available_height
                        image_top = text_box_bottom - Inches(2.50)  # Small padding from bottom
                        
                        # Center the image horizontally if it's smaller than content width
                        if image_width < self.content_width:
                            image_left = self.left_margin + (self.content_width - image_width) / 2
                        else:
                            image_left = self.left_margin
                        
                        # Add image to slide
                        if aspect_ratio < 25:
                            slide.shapes.add_picture(
                                image_path,
                                image_left,
                                image_top,
                                width=image_width,
                                height=image_height
                            )
                        
                    except Exception as e:
                        print(f"Error adding image for question {mcq['number']}: {e}")
                    
        self.add_yellow_border(slide)
                    
    
    def convert_document(self, input_docx, output_pptx):
        """Main conversion function - converts Word to formatted PowerPoint"""
        print(f"Converting: {input_docx} to {output_pptx}")
        
        # Extract images
        images_dict = self.extract_mcq_images(input_docx)
        
        # Extract MCQs
        html_file = self.convert_docx_to_html(input_docx)
        mcqs = self.parse_html(html_file)
        print(f"Found {len(mcqs)} MCQs")

        if not mcqs:
            print("No MCQs found in the document!")
            return False
        
        # Create formatted presentation
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        
        # Create slides for each MCQ
        for i, mcq in enumerate(mcqs):
            # Find the corresponding image for this question number
            mcq_image = None
            for img_dict in images_dict:
                if int(img_dict['name']) == mcq['number']:
                    mcq_image = img_dict
                    break
            
            self.create_formatted_slide(prs, mcq, mcq_image, is_first_slide=(i == 0))
        
        # Save presentation
        prs.save(output_pptx)
        print(f"Formatted presentation saved: {output_pptx}")
        
        return True


def convert_word_to_ppt(input_file, output_file):
    """
    Programmatic interface for converting Word to PowerPoint.
    Can be called from Django views or other Python code.
    
    Args:
        input_file (str): Path to input Word document
        output_file (str): Path for output PowerPoint file
        
    Returns:
        bool: True if conversion successful, False otherwise
    """
    converter = MCQConverter()
    return converter.convert_document(input_file, output_file)
